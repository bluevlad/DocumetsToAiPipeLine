"""Office 문서 변환기.

DOC/DOCX (6,879개), PPT/PPTX (4,856개), XLS/XLSX (5,767개) 처리.
레거시 포맷(.doc, .ppt, .xls)은 바이너리 텍스트 추출 → LibreOffice 폴백.
"""

import asyncio
import logging
import re
from pathlib import Path

import olefile

logger = logging.getLogger(__name__)


class OfficeConverter:
    """MS Office 파일을 텍스트로 변환."""

    SUPPORTED_EXTENSIONS = {
        ".doc", ".docx",
        ".ppt", ".pptx",
        ".xls", ".xlsx",
    }

    def can_handle(self, file_path: Path) -> bool:
        return file_path.suffix.lower() in self.SUPPORTED_EXTENSIONS

    async def convert(self, file_path: Path) -> str:
        """Office 파일을 텍스트로 변환."""
        suffix = file_path.suffix.lower()
        try:
            if suffix == ".docx":
                return await asyncio.to_thread(self._convert_docx, file_path)
            elif suffix == ".doc":
                return await asyncio.to_thread(self._convert_doc_legacy, file_path)
            elif suffix == ".pptx":
                return await asyncio.to_thread(self._convert_pptx, file_path)
            elif suffix == ".ppt":
                return await asyncio.to_thread(self._convert_ppt_legacy, file_path)
            elif suffix == ".xlsx":
                return await asyncio.to_thread(self._convert_xlsx, file_path)
            elif suffix == ".xls":
                return await asyncio.to_thread(self._convert_xls_legacy, file_path)
        except Exception as e:
            logger.warning("Office conversion failed for %s: %s", file_path, e)
        return ""

    # ── Modern formats ──────────────────────────────────────

    def _convert_docx(self, file_path: Path) -> str:
        """DOCX → 텍스트 (python-docx)."""
        from docx import Document

        doc = Document(str(file_path))
        parts = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)

        for table in doc.tables:
            for row in table.rows:
                row_text = "\t".join(
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                )
                if row_text:
                    parts.append(row_text)

        return "\n".join(parts)

    def _convert_pptx(self, file_path: Path) -> str:
        """PPTX → 텍스트 (python-pptx)."""
        from pptx import Presentation

        try:
            prs = Presentation(str(file_path))
        except Exception:
            return ""

        parts = []
        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)

                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = "\t".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            slide_texts.append(row_text)

            if slide_texts:
                parts.append(f"[Slide {slide_idx}]")
                parts.extend(slide_texts)

        return "\n".join(parts)

    def _convert_xlsx(self, file_path: Path) -> str:
        """XLSX → 텍스트 (openpyxl)."""
        from openpyxl import load_workbook

        try:
            wb = load_workbook(str(file_path), read_only=True, data_only=True)
        except Exception:
            return ""

        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_texts = []
            for row in ws.iter_rows(values_only=True):
                row_values = [str(v).strip() for v in row if v is not None]
                if row_values:
                    sheet_texts.append("\t".join(row_values))

            if sheet_texts:
                parts.append(f"[Sheet: {sheet_name}]")
                parts.extend(sheet_texts)

        wb.close()
        return "\n".join(parts)

    # ── Legacy formats (OLE binary extraction) ─────────────

    def _convert_doc_legacy(self, file_path: Path) -> str:
        """레거시 DOC → 텍스트 (OLE WordDocument 스트림에서 추출)."""
        if not olefile.isOleFile(str(file_path)):
            return ""

        ole = olefile.OleFileIO(str(file_path))
        try:
            if not ole.exists("WordDocument"):
                return ""

            # WordDocument 스트림에서 텍스트 추출 시도
            data = ole.openstream("WordDocument").read()
            text = self._extract_text_from_binary(data)

            # 1Word 스트림도 시도
            if not text and ole.exists("1Table"):
                data = ole.openstream("1Table").read()
                text = self._extract_text_from_binary(data)

            return text
        finally:
            ole.close()

    def _convert_ppt_legacy(self, file_path: Path) -> str:
        """레거시 PPT → 텍스트 (OLE에서 텍스트 추출)."""
        if not olefile.isOleFile(str(file_path)):
            return ""

        ole = olefile.OleFileIO(str(file_path))
        try:
            text_parts = []

            # PPT의 텍스트는 PowerPoint Document 또는 Current User 스트림에 있음
            for stream_name in ole.listdir():
                stream_path = "/".join(stream_name)
                try:
                    data = ole.openstream(stream_name).read()
                    extracted = self._extract_text_from_binary(data)
                    if extracted and len(extracted.strip()) > 20:
                        text_parts.append(extracted)
                except Exception:
                    continue

            return "\n".join(text_parts)
        finally:
            ole.close()

    def _convert_xls_legacy(self, file_path: Path) -> str:
        """레거시 XLS → 텍스트 (OLE에서 텍스트 추출)."""
        if not olefile.isOleFile(str(file_path)):
            return ""

        ole = olefile.OleFileIO(str(file_path))
        try:
            text_parts = []
            for stream_name in ole.listdir():
                try:
                    data = ole.openstream(stream_name).read()
                    extracted = self._extract_text_from_binary(data)
                    if extracted and len(extracted.strip()) > 10:
                        text_parts.append(extracted)
                except Exception:
                    continue

            return "\n".join(text_parts)
        finally:
            ole.close()

    def _extract_text_from_binary(self, data: bytes) -> str:
        """바이너리 데이터에서 한글/영문 텍스트 추출."""
        text_parts = []

        # UTF-16LE 텍스트 추출 시도
        try:
            decoded = data.decode("utf-16-le", errors="ignore")
            # 의미 있는 텍스트만 필터링
            clean = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", decoded)
            # 연속된 한글 또는 영문 추출
            matches = re.findall(r"[\uac00-\ud7a3\u3131-\u3163\u3165-\u318e\w\s.,;:!?()\"'-]{5,}", clean)
            if matches:
                text_parts.extend(matches)
        except Exception:
            pass

        # CP949 텍스트 추출 시도
        if not text_parts:
            try:
                decoded = data.decode("cp949", errors="ignore")
                clean = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", decoded)
                matches = re.findall(r"[\uac00-\ud7a3\u3131-\u3163\u3165-\u318e\w\s.,;:!?()\"'-]{5,}", clean)
                if matches:
                    text_parts.extend(matches)
            except Exception:
                pass

        if text_parts:
            result = "\n".join(text_parts)
            # 중복 제거 및 정리
            lines = list(dict.fromkeys(result.split("\n")))
            return "\n".join(line.strip() for line in lines if line.strip())

        return ""
