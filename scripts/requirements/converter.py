"""표 인식 문서 → 텍스트 변환기 (요구사항 추출 전용).

레거시 포맷(.doc/.ppt/.xls)은 LibreOffice(headless)로 최신 포맷으로 변환한 뒤
python-docx / python-pptx / openpyxl 로 **표 구조를 보존**하여 텍스트화한다.
(앱의 OLE 바이너리 정규식 추출보다 표 충실도가 훨씬 높다.)

변환 중간 산출물은 캐시(_converted/)에 보관하여 재실행 시 건너뛴다.
"""
from __future__ import annotations

import subprocess
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

SOFFICE = "/Applications/LibreOffice.app/Contents/MacOS/soffice"

ROOT = Path(__file__).resolve().parents[2]
CONVERTED_DIR = ROOT / "docs" / "requirements" / "_converted"
PROFILE_DIR = CONVERTED_DIR / ".soffice_profile"

LEGACY_TO_MODERN = {".doc": "docx", ".ppt": "pptx", ".xls": "xlsx"}


def _soffice_convert(src: Path, target_ext: str, outdir: Path) -> Path | None:
    """LibreOffice headless 변환. 결과 파일 경로 반환."""
    outdir.mkdir(parents=True, exist_ok=True)
    out = outdir / f"{src.stem}.{target_ext}"
    if out.exists() and out.stat().st_size > 0:
        return out
    try:
        subprocess.run(
            [
                SOFFICE, "--headless",
                f"-env:UserInstallation=file://{PROFILE_DIR}",
                "--convert-to", target_ext, "--outdir", str(outdir), str(src),
            ],
            check=True, capture_output=True, timeout=120,
        )
    except (subprocess.CalledProcessError, subprocess.TimeoutExpired):
        return None
    return out if out.exists() else None


def _docx_to_text(path: Path) -> str:
    doc = Document(str(path))
    parts: list[str] = []
    for p in doc.paragraphs:
        t = p.text.strip()
        if t:
            parts.append(t)
    for ti, table in enumerate(doc.tables):
        parts.append(f"[표 {ti + 1}]")
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            # 병합셀 중복 축약
            dedup, prev = [], None
            for c in cells:
                if c != prev:
                    dedup.append(c)
                prev = c
            line = " | ".join(x for x in dedup if x)
            if line:
                parts.append(line)
    return "\n".join(parts)


def _xlsx_to_text(path: Path) -> str:
    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        rows_text = []
        for row in ws.iter_rows(values_only=True):
            vals = [str(v).strip() for v in row if v is not None and str(v).strip()]
            if vals:
                rows_text.append(" | ".join(vals))
        if rows_text:
            parts.append(f"[시트: {ws.title}]")
            parts.extend(rows_text)
    wb.close()
    return "\n".join(parts)


def _pptx_to_text(path: Path) -> str:
    prs = Presentation(str(path))
    parts: list[str] = []
    for si, slide in enumerate(prs.slides):
        slide_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    slide_parts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    line = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                    if line:
                        slide_parts.append(line)
        if slide_parts:
            parts.append(f"[슬라이드 {si + 1}]")
            parts.extend(slide_parts)
    return "\n".join(parts)


def convert_to_text(abs_path: str) -> str:
    """원본 파일 → 표 보존 텍스트. 실패 시 빈 문자열."""
    src = Path(abs_path)
    ext = src.suffix.lower()
    try:
        if ext == ".txt":
            return src.read_text(encoding="utf-8", errors="ignore")
        if ext == ".docx":
            return _docx_to_text(src)
        if ext == ".xlsx":
            return _xlsx_to_text(src)
        if ext == ".pptx":
            return _pptx_to_text(src)
        if ext in LEGACY_TO_MODERN:
            target = LEGACY_TO_MODERN[ext]
            converted = _soffice_convert(src, target, CONVERTED_DIR)
            if not converted:
                return ""
            if target == "docx":
                return _docx_to_text(converted)
            if target == "xlsx":
                return _xlsx_to_text(converted)
            if target == "pptx":
                return _pptx_to_text(converted)
    except Exception as e:  # noqa: BLE001
        return f""  # 조용히 실패 — 호출측에서 빈 결과 처리
    return ""


if __name__ == "__main__":
    import sys

    txt = convert_to_text(sys.argv[1])
    print(f"[len={len(txt)}]")
    print(txt[:3000])
